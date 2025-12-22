#!/usr/bin/env python3
"""
Create PowerPoint presentation for Legacy Translations Admin System
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Background color
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(26, 54, 93)  # Dark blue
    background.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.33), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.33), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_section_slide(prs, title, color_rgb):
    """Add a section divider slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(*color_rgb)
    background.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12.33), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, items, accent_color=(79, 70, 229)):
    """Add a content slide with bullet points"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # White background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    background.line.fill.background()

    # Accent bar at top
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(0.15))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = RGBColor(*accent_color)
    accent_bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.33), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(30, 30, 30)

    # Content box
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12.33), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = f"• {item}"
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_after = Pt(12)

    return slide

def add_feature_slide(prs, title, features, accent_color=(79, 70, 229)):
    """Add a feature slide with icons/checkmarks"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Light background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(248, 250, 252)
    background.line.fill.background()

    # Accent bar
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(0.15))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = RGBColor(*accent_color)
    accent_bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.33), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(30, 30, 30)

    # Features in two columns
    left_features = features[:len(features)//2 + len(features)%2]
    right_features = features[len(features)//2 + len(features)%2:]

    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(6), Inches(5.5))
    tf = left_box.text_frame
    tf.word_wrap = True

    for i, feat in enumerate(left_features):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"✓ {feat}"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_after = Pt(10)

    # Right column
    if right_features:
        right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.4), Inches(6), Inches(5.5))
        tf = right_box.text_frame
        tf.word_wrap = True

        for i, feat in enumerate(right_features):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"✓ {feat}"
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(50, 50, 50)
            p.space_after = Pt(10)

    return slide

def create_presentation():
    """Create the full presentation"""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # ===== TITLE SLIDE =====
    add_title_slide(prs,
        "Legacy Translations",
        "Sistema de Gestão de Traduções Certificadas")

    # ===== OVERVIEW =====
    add_content_slide(prs, "Visão Geral do Sistema", [
        "Sistema completo de gestão de traduções certificadas",
        "Três níveis de acesso: Admin, Project Manager (PM) e Translator",
        "Fluxo de trabalho integrado: pedido → atribuição → tradução → revisão → entrega",
        "Geração automática de PDFs certificados com letterhead profissional",
        "Notificações por email em todas as etapas do processo",
        "Interface web responsiva e intuitiva"
    ])

    # ===== ADMIN SECTION =====
    add_section_slide(prs, "ADMINISTRADOR", (79, 70, 229))  # Purple

    add_content_slide(prs, "Admin - Gestão de Pedidos", [
        "Visualização de todos os pedidos do sistema",
        "Criação manual de projetos de tradução",
        "Atribuição de Project Managers aos projetos",
        "Atribuição de Tradutores aos projetos",
        "Acompanhamento de status em tempo real",
        "Acesso a valores financeiros (preços e pagamentos)",
        "Exportação de relatórios"
    ], accent_color=(79, 70, 229))

    add_feature_slide(prs, "Admin - Gestão de Usuários", [
        "Criar novos usuários (Admin, PM, Translator, Sales)",
        "Ativar/Desativar contas de usuários",
        "Definir taxas por página para tradutores",
        "Definir taxas por palavra para tradutores",
        "Configurar pares de idiomas dos tradutores",
        "Excluir usuários do sistema",
        "Visualizar lista completa de usuários",
        "Gerenciar permissões de acesso"
    ], accent_color=(79, 70, 229))

    add_content_slide(prs, "Admin - Workspace de Tradução", [
        "Acesso ao workspace de tradução de qualquer projeto",
        "Visualização lado a lado: original e tradução",
        "Seletor de arquivos do projeto",
        "Preview do certificado em tempo real",
        "Geração de PDF final com letterhead",
        "Revisão de traduções finalizadas"
    ], accent_color=(79, 70, 229))

    add_feature_slide(prs, "Admin - Funcionalidades Exclusivas", [
        "Acesso total a preços e valores",
        "Visualização de custos de tradução",
        "Coluna de PM visível na tabela",
        "Criar qualquer tipo de usuário",
        "Configurações do sistema",
        "Relatórios financeiros",
        "Histórico completo de ações",
        "Notificações de sistema"
    ], accent_color=(79, 70, 229))

    # ===== PM SECTION =====
    add_section_slide(prs, "PROJECT MANAGER", (16, 185, 129))  # Green

    add_content_slide(prs, "PM - Gestão de Projetos", [
        "Visualização dos projetos atribuídos ao PM",
        "Atribuição de tradutores aos projetos",
        "Acompanhamento de status das traduções",
        "Coluna de status 'Translation' (Working/Ready/Review)",
        "Edição de atribuições de tradutores",
        "Sem acesso a valores financeiros (preços ocultos)"
    ], accent_color=(16, 185, 129))

    add_feature_slide(prs, "PM - Sistema de Revisão", [
        "Modal de revisão lado a lado",
        "Original à esquerda, tradução à direita",
        "Navegação entre múltiplos documentos",
        "Botão 'Aprovar Tradução'",
        "Botão 'Solicitar Correção'",
        "Campo para comentários de revisão",
        "Notificação automática ao tradutor",
        "Histórico de revisões"
    ], accent_color=(16, 185, 129))

    add_content_slide(prs, "PM - Gestão de Tradutores", [
        "Visualizar lista de tradutores disponíveis",
        "Criar novos tradutores diretamente do modal",
        "Formulário rápido de cadastro de tradutor",
        "Definir nome, email e senha do tradutor",
        "Configurar taxas e pares de idiomas",
        "Atribuição automática após criação"
    ], accent_color=(16, 185, 129))

    add_content_slide(prs, "PM - Workspace de Tradução", [
        "Acesso ao workspace de seus projetos",
        "Seletor de arquivos do projeto",
        "Visualização de documentos originais",
        "Acompanhamento do progresso da tradução",
        "Preview do documento traduzido",
        "Verificação de qualidade"
    ], accent_color=(16, 185, 129))

    # ===== TRANSLATOR SECTION =====
    add_section_slide(prs, "TRADUTOR", (245, 158, 11))  # Orange/Amber

    add_content_slide(prs, "Translator - Meus Projetos", [
        "Dashboard personalizado com projetos atribuídos",
        "Sistema de aceitar/recusar projetos",
        "Visualização de prazos (deadlines)",
        "Status do projeto em tempo real",
        "Notificações de novos projetos",
        "Histórico de projetos concluídos"
    ], accent_color=(245, 158, 11))

    add_feature_slide(prs, "Translator - Workspace de Tradução", [
        "Interface de tradução lado a lado",
        "Documento original à esquerda",
        "Editor de tradução à direita",
        "Seletor de arquivos do projeto",
        "Salvamento automático do progresso",
        "Preview do certificado em tempo real",
        "Geração de PDF com letterhead",
        "Envio da tradução finalizada"
    ], accent_color=(245, 158, 11))

    add_content_slide(prs, "Translator - Fluxo de Trabalho", [
        "1. Receber notificação de novo projeto",
        "2. Aceitar ou recusar o projeto",
        "3. Acessar workspace de tradução",
        "4. Selecionar arquivo para traduzir",
        "5. Realizar tradução com preview em tempo real",
        "6. Gerar PDF certificado",
        "7. Enviar para revisão do PM/Admin"
    ], accent_color=(245, 158, 11))

    add_content_slide(prs, "Translator - Notificações", [
        "Email automático ao receber novo projeto",
        "Detalhes do projeto no email",
        "Link direto para o workspace",
        "Notificação de aprovação/correção",
        "Comentários do revisor visíveis",
        "Sistema de notificações in-app"
    ], accent_color=(245, 158, 11))

    # ===== COMMON FEATURES =====
    add_section_slide(prs, "RECURSOS COMUNS", (59, 130, 246))  # Blue

    add_feature_slide(prs, "PDF Certificado - Recursos", [
        "Letterhead profissional Legacy Translations",
        "Linha azul decorativa sob o header",
        "Seção 'Original Document' com imagem",
        "Seção 'Translation' com texto formatado",
        "Declaração de certificação",
        "Informações do tradutor certificado",
        "Número de certificação único",
        "Data de emissão"
    ], accent_color=(59, 130, 246))

    add_content_slide(prs, "Sistema de Notificações", [
        "Notificações por email em cada etapa",
        "Templates de email profissionais",
        "Notificações in-app em tempo real",
        "Marcação de notificações como lidas",
        "Histórico de notificações",
        "Links diretos para ações"
    ], accent_color=(59, 130, 246))

    add_feature_slide(prs, "Segurança e Autenticação", [
        "Sistema de login seguro",
        "Tokens de autenticação",
        "Controle de acesso por papel (role)",
        "Sessões com expiração",
        "Senhas criptografadas",
        "Logs de atividades",
        "Proteção contra acesso não autorizado",
        "Validação em todas as operações"
    ], accent_color=(59, 130, 246))

    # ===== WORKFLOW SLIDE =====
    add_content_slide(prs, "Fluxo Completo do Sistema", [
        "1. Cliente faz pedido ou Admin cria projeto manualmente",
        "2. Admin/PM atribui tradutor ao projeto",
        "3. Tradutor recebe notificação por email",
        "4. Tradutor aceita o projeto",
        "5. Tradutor realiza a tradução no workspace",
        "6. Tradutor gera PDF e envia para revisão",
        "7. PM/Admin revisa lado a lado",
        "8. Aprovação ou solicitação de correção",
        "9. Entrega final ao cliente"
    ], accent_color=(59, 130, 246))

    # ===== FINAL SLIDE =====
    add_title_slide(prs,
        "Legacy Translations",
        "Sistema Profissional de Gestão de Traduções")

    # Save
    prs.save('/home/user/Legacy-Portal/Legacy_Translations_System.pptx')
    print("Presentation saved: /home/user/Legacy-Portal/Legacy_Translations_System.pptx")

if __name__ == "__main__":
    create_presentation()
