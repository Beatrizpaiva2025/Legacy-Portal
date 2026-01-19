from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor as RgbColor

# Create presentation with 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Legacy Translations Brand Colors (from logo)
NAVY_BLUE = RgbColor(27, 58, 107)      # Dark blue - "LEGACY"
LIGHT_BLUE = RgbColor(91, 155, 213)    # Light blue - "TRANSLATIONS"
DARK_GRAY = RgbColor(51, 51, 51)
LIGHT_GRAY = RgbColor(240, 240, 240)
WHITE = RgbColor(255, 255, 255)

# Logo path
LOGO_PATH = "/home/user/Legacy-Portal/frontend/public/images/legacy-logo.png"

def add_logo(slide, position="header"):
    """Add logo to slide"""
    if position == "header":
        # Small logo in top-left corner
        slide.shapes.add_picture(LOGO_PATH, Inches(0.3), Inches(0.15), height=Inches(0.6))
    elif position == "center":
        # Large centered logo for title slides
        slide.shapes.add_picture(LOGO_PATH, Inches(3.5), Inches(1.5), height=Inches(1.8))
    elif position == "footer":
        # Small logo in bottom corner
        slide.shapes.add_picture(LOGO_PATH, Inches(0.3), Inches(6.7), height=Inches(0.5))

def add_title_slide(prs, title, subtitle=""):
    """Add a title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Background - white with navy accent bar at bottom
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(6.5), prs.slide_width, Inches(1))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = NAVY_BLUE
    accent_bar.line.fill.background()

    # Logo - large centered
    add_logo(slide, "center")

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = NAVY_BLUE
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        p = tf.add_paragraph()
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = LIGHT_BLUE
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_section_slide(prs, title, section_number=""):
    """Add a section divider slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Navy accent bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3), prs.slide_width, Inches(1.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY_BLUE
    shape.line.fill.background()

    # Light blue accent line
    accent_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.5), prs.slide_width, Inches(0.1))
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = LIGHT_BLUE
    accent_line.line.fill.background()

    # Section number
    if section_number:
        num_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(2), Inches(1))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = section_number
        p.font.size = Pt(72)
        p.font.bold = True
        p.font.color.rgb = LIGHT_BLUE

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.15), Inches(12.333), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Logo in corner
    add_logo(slide, "footer")

    return slide

def add_content_slide(prs, title, content_lines, has_image_placeholder=True):
    """Add a content slide with bullet points and image placeholder"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title bar - Navy blue
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = NAVY_BLUE
    title_bar.line.fill.background()

    # Light blue accent line under title
    accent_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.2), prs.slide_width, Inches(0.05))
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = LIGHT_BLUE
    accent_line.line.fill.background()

    # Logo in title bar (right side)
    slide.shapes.add_picture(LOGO_PATH, Inches(11.5), Inches(0.25), height=Inches(0.7))

    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(10.5), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Content area
    if has_image_placeholder:
        # Left side - text content
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(5.5), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True

        for i, line in enumerate(content_lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(18)
            p.font.color.rgb = DARK_GRAY
            p.space_after = Pt(12)

        # Right side - image placeholder with light blue border
        img_placeholder = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(6.5), Inches(1.6),
            Inches(6.333), Inches(5.4)
        )
        img_placeholder.fill.solid()
        img_placeholder.fill.fore_color.rgb = LIGHT_GRAY
        img_placeholder.line.color.rgb = LIGHT_BLUE
        img_placeholder.line.width = Pt(3)

        # Placeholder text
        img_text = slide.shapes.add_textbox(Inches(6.5), Inches(3.9), Inches(6.333), Inches(1))
        tf = img_text.text_frame
        p = tf.paragraphs[0]
        p.text = "INSERT SCREENSHOT HERE"
        p.font.size = Pt(18)
        p.font.color.rgb = LIGHT_BLUE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
    else:
        # Full width content
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(12.333), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True

        for i, line in enumerate(content_lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(20)
            p.font.color.rgb = DARK_GRAY
            p.space_after = Pt(14)

    return slide

def add_table_slide(prs, title, headers, rows):
    """Add a slide with a table"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title bar - Navy blue
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = NAVY_BLUE
    title_bar.line.fill.background()

    # Light blue accent line
    accent_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.2), prs.slide_width, Inches(0.05))
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = LIGHT_BLUE
    accent_line.line.fill.background()

    # Logo in title bar
    slide.shapes.add_picture(LOGO_PATH, Inches(11.5), Inches(0.25), height=Inches(0.7))

    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(10.5), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Table
    cols = len(headers)
    table_rows = len(rows) + 1  # +1 for header

    table = slide.shapes.add_table(table_rows, cols, Inches(0.5), Inches(1.6), Inches(12.333), Inches(5)).table

    # Set column widths
    col_width = Inches(12.333 / cols)
    for i in range(cols):
        table.columns[i].width = col_width

    # Header row - Navy blue background
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY_BLUE
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    # Data rows
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = cell_text
            # Alternate row colors
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RgbColor(240, 248, 255)  # Light blue tint
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.color.rgb = DARK_GRAY
            p.alignment = PP_ALIGN.CENTER

    return slide

# ==================== CREATE SLIDES ====================

# Slide 1: Title
add_title_slide(prs, "Partner Portal Guide", "Step-by-Step Instructions")

# Slide 2: Table of Contents
add_content_slide(prs, "Table of Contents", [
    "1. Partner Registration",
    "2. Login to the Portal",
    "3. Creating a New Order",
    "4. Tracking Orders (My Orders)",
    "5. Invoices and Payments",
    "6. Messages",
    "7. Payment Plan",
    "8. Support and Help"
], has_image_placeholder=False)

# ==================== SECTION 1: REGISTRATION ====================
add_section_slide(prs, "Partner Registration", "01")

add_content_slide(prs, "Step 1: Access the Registration Page", [
    "1. Go to the Partner Portal:",
    "   https://legacy-portal-frontend.onrender.com/#/partner",
    "",
    "2. On the login screen, click:",
    '   "Don\'t have an account? Register here"',
    "",
    "3. The registration form will appear."
])

add_content_slide(prs, "Step 2: Fill in Company Information", [
    "Required fields:",
    "   Company Name - Your company's legal name",
    "   Contact Name - Primary contact person",
    "   Company Email - Used for login and notifications",
    "   Password - Minimum 6 characters",
    "",
    "Optional fields:",
    "   Phone - Contact phone number",
    "   Company Address"
])

add_content_slide(prs, "Step 3: Select Monthly Volume", [
    "Choose your estimated monthly translation volume:",
    "",
    "   1-10 pages per month",
    "   11-50 pages per month",
    "   51-100 pages per month",
    "   100+ pages per month",
    "",
    "This helps us better serve your needs."
])

add_content_slide(prs, "Step 4: Complete Registration", [
    "1. Review your information",
    "",
    '2. Check the box: "I agree to the Partner Agreement"',
    "",
    '3. Click the green "Create Account" button',
    "",
    "4. You'll receive a confirmation email",
    "",
    "5. Your account is now active!"
])

# ==================== SECTION 2: LOGIN ====================
add_section_slide(prs, "Login to the Portal", "02")

add_content_slide(prs, "Login Steps", [
    "1. Go to the Partner Portal URL",
    "",
    "2. Enter your Company Email",
    "",
    "3. Enter your Password",
    "",
    '4. Click the "Login" button',
    "",
    "You'll be redirected to the dashboard."
])

add_content_slide(prs, "Forgot Password?", [
    '1. Click "Forgot password?" on the login page',
    "",
    "2. Enter your registered email address",
    "",
    '3. Click "Send Reset Link"',
    "",
    "4. Check your email for the reset link",
    "",
    "5. Click the link and set a new password"
])

# ==================== SECTION 3: NEW ORDER ====================
add_section_slide(prs, "Creating a New Order", "03")

add_content_slide(prs, "Step 1: Access New Order Page", [
    "After logging in, you'll be on the",
    "New Order page automatically.",
    "",
    'If not, click "+ New Order" in the',
    "left sidebar menu.",
    "",
    "The order form will appear."
])

add_content_slide(prs, "Step 2: Client Information", [
    "Fill in the Client Information section:",
    "",
    "   Client Name - End client's full name",
    "",
    "   Client Email - Client will receive",
    "   status notifications at this email"
])

add_table_slide(prs, "Step 3: Select Service Type",
    ["Service", "Price/Page", "When to Use"],
    [
        ["Certified Translation", "$24.99", "USCIS, legal, immigration"],
        ["Standard Translation", "$19.99", "General use, no certification"],
        ["Sworn Translation", "$55.00", "International use outside USA"],
        ["RMV Certified", "$24.99", "Massachusetts DMV/RMV"]
    ]
)

add_content_slide(prs, "Step 4: Select Languages", [
    "Choose the language pair:",
    "",
    "   Translate From: Select the original",
    "   document language (e.g., Portuguese)",
    "",
    "   Translate To: Select the target",
    "   language (e.g., English)",
    "",
    "55+ languages available!"
])

add_content_slide(prs, "Step 5: Upload Documents", [
    "In the Upload Documents section:",
    "",
    "   Drag and drop files into the area, OR",
    "   Click to browse and select files",
    "",
    "Accepted formats:",
    "PDF, DOCX, TXT, JPEG, PNG, TIFF",
    "",
    "Pages are calculated automatically",
    "(250 words = 1 page)"
])

add_table_slide(prs, "Step 6: Select Urgency (Optional)",
    ["Option", "Delivery Time", "Additional Fee"],
    [
        ["Standard", "2-3 business days", "No fee"],
        ["Priority", "24 hours", "+25%"],
        ["Urgent", "12 hours", "+100%"]
    ]
)

add_content_slide(prs, "Step 7: Physical Copy (Optional)", [
    'Check "Physical copy required" if the',
    "client needs a mailed hard copy.",
    "",
    "Cost: $18.99 (USPS Priority Mail)",
    "Delivery: 1-3 business days",
    "",
    "Fill in the shipping address:",
    "   Street Address",
    "   City, State, ZIP Code",
    "",
    "Note: Required for RMV Certified orders"
])

add_content_slide(prs, "Step 8: Review Quote Summary", [
    "Check the Quote Summary panel on the right:",
    "",
    "   Service type selected",
    "   Total pages",
    "   Base price",
    "   Urgency fee (if applicable)",
    "   Shipping fee (if applicable)",
    "   Coupon discount (if applied)",
    "",
    "   TOTAL PRICE"
])

add_content_slide(prs, "Step 9: Apply Coupon (Optional)", [
    "If you have a coupon code:",
    "",
    '1. Enter the code in the "Enter code" field',
    "",
    '2. Click the "Apply" button',
    "",
    "3. The discount will be reflected",
    "   in the Quote Summary"
])

add_content_slide(prs, "Step 10: Select Payment Method", [
    "Choose your payment method:",
    "",
    "OPTION 1: Invoice",
    "   Receive invoice by email",
    "   Pay later via card or Zelle",
    "",
    "OPTION 2: Zelle (Direct Payment)",
    "   Send payment to: contact@legacytranslations.com",
    "   Include company name in memo",
    "   Upload payment receipt"
])

add_content_slide(prs, "Step 11: Submit Order", [
    '1. Click the green "Create Order" button',
    "",
    "2. Wait for confirmation message",
    "",
    "3. You'll receive an email with order details",
    "",
    "4. Track your order in 'My Orders'",
    "",
    "Done! Your order is being processed."
])

# ==================== SECTION 4: MY ORDERS ====================
add_section_slide(prs, "Tracking Orders", "04")

add_content_slide(prs, "Access My Orders", [
    'Click "My Orders" in the left sidebar.',
    "",
    "You'll see all your orders with:",
    "   Order number",
    "   Client name",
    "   Service type",
    "   Status (Paid/Pending/Overdue)",
    "   Translation stage",
    "   Amount"
])

add_content_slide(prs, "Order Status Badges", [
    "Status indicators:",
    "",
    "   PAID (green) - Payment received",
    "",
    "   PENDING (yellow) - Awaiting payment",
    "",
    "   OVERDUE (red) - Payment past due"
])

add_content_slide(prs, "Translation Stages", [
    "Track your order's progress:",
    "",
    "   Received - Order received",
    "",
    "   Processing - Translation in progress",
    "",
    "   Completed - Translation finished",
    "",
    "   Delivered - Sent to client"
])

add_content_slide(prs, "Filter Orders", [
    "Use the filter buttons at the top:",
    "",
    "   All - View all orders",
    "   Pending - Only pending orders",
    "   Paid - Only paid orders",
    "   Overdue - Only overdue orders",
    "",
    "Click an order to expand and see details."
])

add_content_slide(prs, "Download Translated Documents", [
    "When translation is complete:",
    "",
    "1. Click on the order to expand it",
    "",
    "2. Find the Documents section",
    "",
    '3. Click "Download" next to the',
    "   translated file",
    "",
    "4. The file will download to your device"
])

# ==================== SECTION 5: INVOICES ====================
add_section_slide(prs, "Invoices and Payments", "05")

add_content_slide(prs, "Access Invoices", [
    'Click "Invoices" in the left sidebar.',
    "",
    "You'll see all your invoices with:",
    "   Invoice number",
    "   Issue date",
    "   Due date",
    "   Status (Paid/Pending/Overdue)",
    "   Total amount"
])

add_content_slide(prs, "View Invoice Details", [
    '1. Click "View Details" on any invoice',
    "",
    "2. A modal will show:",
    "   Invoice number and dates",
    "   List of orders included",
    "   Itemized amounts",
    "   Payment history"
])

add_content_slide(prs, "Pay with Credit Card", [
    '1. Click "Pay Now" on the invoice',
    "",
    '2. Select "Pay with Card"',
    "",
    "3. You'll be redirected to Stripe",
    "",
    "4. Enter your card details",
    "",
    "5. Confirm payment",
    "",
    "6. Invoice marked as Paid automatically"
])

add_content_slide(prs, "Pay with Zelle", [
    '1. Click "Pay Now" on the invoice',
    "",
    '2. Select "Pay with Zelle"',
    "",
    "3. Send payment to:",
    "   payments@legacytranslations.com",
    "",
    "4. Include company name + invoice #",
    "",
    "5. Upload your payment receipt",
    "",
    '6. Click "Submit Receipt"'
])

# ==================== SECTION 6: MESSAGES ====================
add_section_slide(prs, "Messages", "06")

add_content_slide(prs, "Access Messages", [
    'Click "Messages" in the left sidebar.',
    "",
    "Two tabs available:",
    "",
    "   Conversations - Message threads",
    "   grouped by order",
    "",
    "   Notifications - System updates",
    "   about orders and payments"
])

add_content_slide(prs, "View Conversations", [
    'In the "Conversations" tab:',
    "",
    "   See all message threads",
    "   Unread messages shown in bold",
    "   Click a conversation to expand",
    "   View full message history"
])

add_content_slide(prs, "Send a Message", [
    'To send a message about an order:',
    "",
    '1. Go to "My Orders"',
    "",
    "2. Click on the order",
    "",
    '3. Click "Send Message" button',
    "",
    "4. Type your message",
    "",
    '5. Click "Send"'
])

# ==================== SECTION 7: PAYMENT PLAN ====================
add_section_slide(prs, "Payment Plan", "07")

add_content_slide(prs, "Access Payment Plan", [
    'Click "Payment Plan" in the left sidebar.',
    "",
    "Here you can:",
    "   View your current plan",
    "   Check qualification progress",
    "   Request plan upgrades"
])

add_table_slide(prs, "Available Payment Plans",
    ["Plan", "Billing", "Requirement"],
    [
        ["Pay Per Order", "Per order", "Default for new partners"],
        ["Biweekly Invoice", "Every 2 weeks", "3 paid orders"],
        ["Monthly Invoice", "Monthly", "3 paid orders"]
    ]
)

add_content_slide(prs, "Qualify for Invoice Plans", [
    "To unlock invoice payment plans:",
    "",
    "   Complete 3 paid orders",
    "",
    "Your progress is shown on the",
    "Payment Plan page.",
    "",
    "Once qualified, you can request",
    "an upgrade to invoice billing."
])

add_content_slide(prs, "Request Plan Upgrade", [
    "After completing 3 orders:",
    "",
    "1. Go to Payment Plan page",
    "",
    "2. Select your preferred plan:",
    "   Biweekly Invoice, OR",
    "   Monthly Invoice",
    "",
    '3. Click "Request Upgrade"',
    "",
    "4. Wait for admin approval"
])

# ==================== SECTION 8: SUPPORT ====================
add_section_slide(prs, "Support and Help", "08")

add_content_slide(prs, "Support Options", [
    "Option 1: Live Chat",
    "Click the chat icon in the",
    "bottom-right corner",
    "",
    "Option 2: Support Form",
    'Click "Having trouble uploading?"',
    "on the New Order page",
    "",
    "Option 3: Email",
    "contact@legacytranslations.com"
])

add_content_slide(prs, "Change Portal Language", [
    "The portal is available in:",
    "",
    "   English (US)",
    "   Spanish (ES)",
    "   Portuguese (BR)",
    "",
    "To change language:",
    "Click the flag icons in the",
    "top-right corner of the screen"
])

# ==================== SUMMARY ====================
add_section_slide(prs, "Quick Reference", "")

add_table_slide(prs, "Navigation Summary",
    ["Action", "Where to Find"],
    [
        ["Create Order", "Sidebar > New Order"],
        ["View Orders", "Sidebar > My Orders"],
        ["Pay Invoices", "Sidebar > Invoices"],
        ["Messages", "Sidebar > Messages"],
        ["Payment Plan", "Sidebar > Payment Plan"],
        ["Change Language", "Flags in top-right"],
        ["Get Support", "Chat icon bottom-right"],
        ["Logout", "Logout button in sidebar"]
    ]
)

# Final slide
add_title_slide(prs, "Questions?", "Contact us at contact@legacytranslations.com")

# Save presentation
output_path = '/home/user/Legacy-Portal/Partner_Portal_Guide.pptx'
prs.save(output_path)
print(f"Presentation saved: {output_path}")
print(f"Total slides: {len(prs.slides)}")
